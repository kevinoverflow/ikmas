"""PPTX extractor for IKMAS ingestion system."""

from typing import List
from langchain_core.documents import Document
from app.rag.storage import StoredFile
from pptx import Presentation


def extract_pptx_documents(stored: StoredFile) -> List[Document]:
    """
    Extract documents from a PPTX file.
    
    Args:
        stored: StoredFile object containing the PPTX file
        
    Returns:
        List of Document objects with metadata
    """
    # Read the PPTX file content
    presentation = Presentation(stored.path)
    
    # Extract text from slides
    all_text = []
    for slide_num, slide in enumerate(presentation.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_text.append(text)
        
        if slide_text:
            all_text.append(f"Slide {slide_num + 1}:\n" + "\n".join(slide_text))
    
    content = '\n\n'.join(all_text)
    
    # Create a single document with metadata
    doc = Document(
        page_content=content,
        metadata={
            "file_id": stored.file_id,
            "source": stored.original_name,
            "chunk_type": "pptx_text"
        }
    )
    
    return [doc]